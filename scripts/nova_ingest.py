#!/usr/bin/env python3
"""
nova_ingest.py — Digest files into Nova's vector memory.

Supports: PDF, DOCX, XLSX, PPTX, TXT, MD, CSV, RTF, and plain text.
Can read from a local file path OR download from a Slack file URL.

Usage:
  nova_ingest.py /path/to/file.pdf [--topic "car manual"] [--source "jordan"]
  nova_ingest.py --slack-url <url> --filename "manual.pdf" [--topic "..."]
  nova_ingest.py --slack-file-id <file_id> [--topic "..."]

Written by Jordan Koch.
"""

import argparse
import json
import os
import subprocess
import sys
import tempfile
import urllib.request
from pathlib import Path

SCRIPTS    = Path.home() / ".openclaw/scripts"
VECTOR_URL = "http://127.0.0.1:18790/remember"
CHUNK_SIZE = 800   # characters per memory chunk (fits in embedding context)
CHUNK_OVERLAP = 100

sys.path.insert(0, str(SCRIPTS))
import nova_config


# ── Text Extractors ───────────────────────────────────────────────────────────

def extract_pdf(path: str) -> str:
    """Extract text from PDF using pdfplumber (best quality) with pdftotext fallback."""
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
        return "\n\n".join(text_parts)
    except Exception:
        pass
    # Fallback: pdftotext
    try:
        result = subprocess.run(
            ["pdftotext", path, "-"],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return result.stdout
    except Exception:
        pass
    return ""


def extract_docx(path: str) -> str:
    """Extract text from Word documents."""
    try:
        import docx
        doc = docx.Document(path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"[DOCX extract error: {e}]"


def extract_xlsx(path: str) -> str:
    """Extract text from Excel spreadsheets."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"[XLSX extract error: {e}]"


def extract_pptx(path: str) -> str:
    """Extract text from PowerPoint presentations."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if slide_parts:
                parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_parts))
        return "\n\n".join(parts)
    except Exception as e:
        return f"[PPTX extract error: {e}]"


def extract_rtf(path: str) -> str:
    """Extract text from RTF using macOS textutil."""
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", path],
            capture_output=True, text=True, timeout=15
        )
        return result.stdout if result.returncode == 0 else ""
    except Exception:
        return ""


def extract_text(path: str, suffix: str) -> str:
    """Route to the right extractor based on file extension."""
    suffix = suffix.lower().lstrip(".")
    if suffix == "pdf":
        return extract_pdf(path)
    elif suffix in ("docx", "doc"):
        if suffix == "doc":
            # Convert .doc to .docx via textutil first
            with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
                tmp_path = tmp.name
            subprocess.run(["textutil", "-convert", "docx", "-output", tmp_path, path],
                           capture_output=True, timeout=15)
            text = extract_docx(tmp_path)
            os.unlink(tmp_path)
            return text
        return extract_docx(path)
    elif suffix in ("xlsx", "xls"):
        return extract_xlsx(path)
    elif suffix in ("pptx", "ppt"):
        return extract_pptx(path)
    elif suffix == "rtf":
        return extract_rtf(path)
    elif suffix in ("txt", "md", "csv", "json", "yaml", "yml", "log", "py",
                    "js", "ts", "swift", "sh", "html", "xml"):
        return Path(path).read_text(encoding="utf-8", errors="ignore")
    else:
        # Try plain text as fallback
        try:
            return Path(path).read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return ""


# ── Chunking ──────────────────────────────────────────────────────────────────

def chunk_text(text: str, filename: str) -> list[str]:
    """Split text into overlapping chunks suitable for embedding."""
    text = text.strip()
    if not text:
        return []

    # Split on paragraph boundaries first
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

    chunks = []
    current = []
    current_len = 0

    for para in paragraphs:
        para_len = len(para)
        if current_len + para_len > CHUNK_SIZE and current:
            chunk_text = "\n\n".join(current)
            chunks.append(f"[From: {filename}]\n{chunk_text}")
            # Keep last paragraph for overlap
            overlap = current[-1] if current else ""
            current = [overlap] if overlap else []
            current_len = len(overlap)
        current.append(para)
        current_len += para_len

    if current:
        chunk_text = "\n\n".join(current)
        chunks.append(f"[From: {filename}]\n{chunk_text}")

    return chunks


# ── Vector Memory ─────────────────────────────────────────────────────────────

def store_chunks(chunks: list[str], source: str, topic: str = "") -> int:
    """Store text chunks in vector memory. Returns number stored."""
    metadata = {"source_type": source}
    if topic:
        metadata["topic"] = topic

    stored = 0
    for chunk in chunks:
        try:
            payload = json.dumps({
                "text": chunk,
                "source": source,
                "metadata": metadata
            }).encode()
            req = urllib.request.Request(
                VECTOR_URL, data=payload,
                headers={"Content-Type": "application/json"}, method="POST"
            )
            with urllib.request.urlopen(req, timeout=15):
                stored += 1
        except Exception as e:
            print(f"[nova_ingest] store error: {e}", file=sys.stderr)

    return stored


# ── Slack Download ────────────────────────────────────────────────────────────

def download_slack_file(file_id: str = None, url: str = None,
                         filename: str = "slack_file") -> tuple[str, str]:
    """Download a file from Slack. Returns (local_path, filename)."""
    token = nova_config.slack_bot_token()
    if not token:
        raise ValueError("Slack token unavailable (Keychain locked?)")

    if file_id and not url:
        # Get file info to find download URL
        # Requires files:read scope on the bot token.
        # If missing, add it at api.slack.com → OAuth & Permissions → Bot Token Scopes
        req = urllib.request.Request(
            f"https://slack.com/api/files.info?file={file_id}",
            headers={"Authorization": f"Bearer {token}"}
        )
        with urllib.request.urlopen(req, timeout=15) as r:
            info = json.loads(r.read())
        if not info.get("ok"):
            error = info.get("error", "unknown")
            if error == "missing_scope":
                raise ValueError(
                    "Bot token missing 'files:read' scope.\n"
                    "Fix: go to api.slack.com → Your Apps → Nova → "
                    "OAuth & Permissions → Bot Token Scopes → add files:read → Reinstall App"
                )
            raise ValueError(f"files.info failed: {error}")
        file_info = info["file"]
        url = file_info.get("url_private_download") or file_info.get("url_private")
        filename = file_info.get("name", filename)

    # Download the file
    req = urllib.request.Request(url, headers={"Authorization": f"Bearer {token}"})
    suffix = Path(filename).suffix or ".tmp"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp_path = tmp.name
        with urllib.request.urlopen(req, timeout=60) as r:
            content = r.read()
        tmp.write(content)

    # Sanity check — Slack redirects to HTML auth page if scope missing
    if content[:15].lower().startswith(b"<!doctype html>") or content[:6] == b"<html>":
        os.unlink(tmp_path)
        raise ValueError(
            "Downloaded an HTML page instead of the file — token missing 'files:read' scope.\n"
            "Fix: api.slack.com → Your Apps → Nova → OAuth & Permissions → add files:read → Reinstall"
        )

    return tmp_path, filename


# ── Main ──────────────────────────────────────────────────────────────────────

def ingest(file_path: str, filename: str, topic: str = "", source: str = "document") -> dict:
    """Ingest a file and return summary dict."""
    suffix = Path(filename).suffix

    print(f"[nova_ingest] Extracting text from {filename}...", flush=True)
    text = extract_text(file_path, suffix)

    if not text.strip():
        return {"ok": False, "error": "No text extracted", "filename": filename}

    word_count = len(text.split())
    chunks = chunk_text(text, filename)

    print(f"[nova_ingest] {word_count} words → {len(chunks)} chunks. Storing...", flush=True)
    stored = store_chunks(chunks, source=source, topic=topic or Path(filename).stem)

    return {
        "ok": True,
        "filename": filename,
        "words": word_count,
        "chunks": len(chunks),
        "stored": stored,
        "topic": topic or Path(filename).stem,
        "preview": text[:300]
    }


def main():
    parser = argparse.ArgumentParser(description="Ingest a file into Nova's vector memory")
    parser.add_argument("file", nargs="?", help="Local file path")
    parser.add_argument("--slack-file-id", help="Slack file ID to download")
    parser.add_argument("--slack-url", help="Slack private download URL")
    parser.add_argument("--filename", help="Filename hint (for Slack downloads)")
    parser.add_argument("--topic", default="", help="Topic label for memory")
    parser.add_argument("--source", default="document", help="Source label")
    args = parser.parse_args()

    tmp_to_delete = None

    if args.slack_file_id or args.slack_url:
        print(f"[nova_ingest] Downloading from Slack...", flush=True)
        file_path, filename = download_slack_file(
            file_id=args.slack_file_id,
            url=args.slack_url,
            filename=args.filename or "slack_file"
        )
        tmp_to_delete = file_path
    elif args.file:
        file_path = args.file
        filename = args.filename or Path(file_path).name
    else:
        parser.print_help()
        sys.exit(1)

    try:
        result = ingest(file_path, filename, topic=args.topic, source=args.source)
        print(json.dumps(result, indent=2))
        if result["ok"]:
            print(f"\n✅ Stored {result['stored']} chunks from \"{result['filename']}\" "
                  f"({result['words']} words) under topic \"{result['topic']}\"")
        else:
            print(f"\n❌ Failed: {result.get('error')}")
            sys.exit(1)
    finally:
        if tmp_to_delete and os.path.exists(tmp_to_delete):
            os.unlink(tmp_to_delete)


if __name__ == "__main__":
    main()
