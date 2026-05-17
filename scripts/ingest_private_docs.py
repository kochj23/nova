#!/usr/bin/env python3
"""
ingest_private_docs.py — Ingest work documents into Nova's PRIVATE memory.

SECURITY: All memories marked source="private_document" with metadata.private=true.
These MUST never be:
  - Included in Slack/Discord/Signal messages
  - Sent to OpenRouter or any cloud LLM
  - Included in any automated report or digest
  - Surfaced unless Jordan explicitly asks Nova for them in a private context

Source: /Volumes/nas/GoogleDriveBackups/Docs/
Written by Jordan Koch.
"""
import json, os, sys, time, subprocess, urllib.request
from pathlib import Path
from datetime import datetime

MEMORY_URL = "http://192.168.1.6:18790/remember"
DOCS_DIR = Path("/Volumes/nas/GoogleDriveBackups/Docs")
LOG_FILE = Path.home() / ".openclaw/logs/ingest-private-docs.log"
CHUNK_SIZE = 400  # words per chunk (smaller for denser recall)

count = 0
failed = 0
skipped = 0
files_processed = 0

def log(msg):
    ts = datetime.now().strftime("%H:%M:%S")
    line = f"[ingest_private {ts}] {msg}"
    print(line, flush=True)
    with open(LOG_FILE, "a") as f:
        f.write(line + "\n")

def remember(text, filename, filepath_rel):
    global count, failed
    if not text.strip() or len(text.strip()) < 30:
        return False
    payload = json.dumps({
        "text": text[:2000],
        "source": "private_document",
        "metadata": {
            "type": "work_document",
            "private": True,
            "file": filename,
            "path": filepath_rel,
            "ingested": datetime.now().isoformat()
        }
    }).encode()
    req = urllib.request.Request(MEMORY_URL, data=payload, headers={"Content-Type": "application/json"}, method="POST")
    try:
        with urllib.request.urlopen(req, timeout=15):
            count += 1
            return True
    except:
        failed += 1
        if failed % 50 == 0:
            time.sleep(2)
        return False

def extract_pdf(filepath):
    try:
        result = subprocess.run(
            ["pdftotext", "-layout", str(filepath), "-"],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode == 0:
            return result.stdout
    except FileNotFoundError:
        pass
    try:
        import PyPDF2
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            return "\n".join(p.extract_text() or "" for p in reader.pages[:80])
    except:
        pass
    return ""

def extract_docx(filepath):
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(filepath)],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return result.stdout
    except:
        pass
    try:
        import docx
        doc = docx.Document(str(filepath))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except:
        pass
    return ""

def extract_doc(filepath):
    """Extract from legacy .doc format using textutil (macOS)."""
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(filepath)],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return result.stdout
    except:
        pass
    return ""

def extract_xlsx(filepath):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
        text_parts = []
        for sheet in wb.worksheets[:10]:
            for row in sheet.iter_rows(max_row=500, values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    text_parts.append(" | ".join(cells))
        return "\n".join(text_parts)
    except:
        return ""

def extract_xls(filepath):
    """Try to extract from legacy .xls format."""
    try:
        import xlrd
        wb = xlrd.open_workbook(str(filepath))
        text_parts = []
        for sheet in wb.sheets()[:10]:
            for row_idx in range(min(sheet.nrows, 500)):
                cells = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols) 
                         if str(sheet.cell_value(row_idx, col)).strip()]
                if cells:
                    text_parts.append(" | ".join(cells))
        return "\n".join(text_parts)
    except:
        return ""

def extract_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(str(filepath))
        text_parts = []
        for slide in prs.slides[:50]:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    except:
        return ""

def chunk_text(text, chunk_size=CHUNK_SIZE):
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size):
        chunk = " ".join(words[i:i + chunk_size])
        if len(chunk.strip()) > 50:
            chunks.append(chunk)
    return chunks

def process_file(filepath):
    global skipped, files_processed
    ext = filepath.suffix.lower().rstrip("*")
    
    skip_exts = {'.ds_store', '.exe', '.zip', '.rar', '.jpg', '.png', '.gif', 
                 '.bmp', '.ico', '.dll', '.vsd', '.msi', '.tmp', '.log'}
    if ext in skip_exts:
        skipped += 1
        return 0
    
    if ext == '.pdf':
        text = extract_pdf(filepath)
    elif ext == '.docx':
        text = extract_docx(filepath)
    elif ext == '.doc':
        text = extract_doc(filepath)
    elif ext in ('.xlsx',):
        text = extract_xlsx(filepath)
    elif ext in ('.xls',):
        text = extract_xls(filepath)
    elif ext in ('.pptx', '.ppt'):
        text = extract_pptx(filepath)
    elif ext in ('.txt', '.md', '.csv', '.xml', '.html', '.htm'):
        try:
            text = filepath.read_text(encoding="utf-8", errors="ignore")[:50000]
        except:
            text = ""
    else:
        skipped += 1
        return 0
    
    if not text or len(text.strip()) < 50:
        skipped += 1
        return 0
    
    files_processed += 1
    rel_path = filepath.name
    chunks = chunk_text(text)
    ingested = 0
    
    for chunk in chunks:
        if remember(chunk, filepath.name, rel_path):
            ingested += 1
        if ingested % 100 == 0 and ingested > 0:
            time.sleep(0.5)
    
    return ingested

def main():
    global count, failed, skipped, files_processed
    log("=" * 60)
    log("PRIVATE DOCUMENT INGESTION")
    log(f"Source: {DOCS_DIR}")
    log(f"Security: source=private_document, metadata.private=true")
    log("=" * 60)

    all_files = sorted([f for f in DOCS_DIR.iterdir() if f.is_file() and not f.name.startswith(".")])
    log(f"Found {len(all_files)} files")
    
    for i, filepath in enumerate(all_files):
        if (i + 1) % 100 == 0:
            log(f"Progress: {i+1}/{len(all_files)} files | {count:,} chunks ingested | {failed} failed | {skipped} skipped")
            time.sleep(1)
        
        try:
            ingested = process_file(filepath)
        except Exception as e:
            log(f"  Error processing {filepath.name}: {e}")
            skipped += 1
    
    log("=" * 60)
    log(f"PRIVATE INGESTION COMPLETE")
    log(f"  Total files: {len(all_files)}")
    log(f"  Files with content: {files_processed}")
    log(f"  Chunks ingested: {count:,}")
    log(f"  Failed: {failed}")
    log(f"  Skipped: {skipped}")
    log("=" * 60)
    
    # Post summary (counts only, no content) to Slack
    sys.path.insert(0, str(Path.home() / ".openclaw/scripts"))
    import nova_config
    nova_config.post_both(
        f":lock: *Private Document Ingestion Complete*\n"
        f"  Files: {files_processed} processed ({skipped} skipped)\n"
        f"  Chunks: {count:,} ingested\n"
        f"  Source: `private_document` (never shared externally)",
        slack_channel=nova_config.SLACK_NOTIFY
    )

if __name__ == "__main__":
    main()
